"""
lol-extract — an Open WebUI "External Document Loader" for the LlmOnLan farm.

One shared document-extraction service on the farm box. Each client's OWUI is
configured (by the shell's config-bridge, from the beacon) with:

    CONTENT_EXTRACTION_ENGINE = external
    EXTERNAL_DOCUMENT_LOADER_URL = http://<farm>:<port>
    EXTERNAL_DOCUMENT_LOADER_API_KEY = <key>

so on every uploaded file OWUI does (verified against open-webui v0.10.2,
backend/open_webui/retrieval/loaders/external_document.py):

    PUT  {URL}/process        body = raw file bytes
    headers: Content-Type: <mime>, Authorization: Bearer <key>, X-Filename: <urlquoted name>

and expects back a JSON list of {"page_content": str, "metadata": {...}} (one per
page). This is the ONLY way an external service ever receives an uploaded file in
OWUI 0.10.2 — external OpenAPI *tool servers* never get the bytes — so both the
"searchable scanned docs (RAG)" and the "vision-model OCR transcript" goals funnel
through this single engine, which routes internally:

  • images + scanned / image-only PDF pages  -> Ollama-OCR (a vision model on the
    farm's LOCAL Ollama /api/generate) -> structured markdown.
  • born-digital PDFs / docx / plaintext / html -> fast local text extraction.
  • (optional, OCR_DOCLING=1) Docling in-process for richer office/table fidelity.

Env (set by farm/src/extract.js spawnExtract):
  EXTRACT_API_KEY  bearer token OWUI must send (matches the beacon-advertised key)
  EXTRACT_PORT     bind port (informational; uvicorn is launched with --port)
  OCR_MODEL        vision model tag for Ollama-OCR (e.g. gemma4:12b)
  OCR_OLLAMA_URL   the local Ollama NATIVE generate endpoint (…:11434/api/generate)
  OCR_FORMAT       markdown|text|json|structured|key_value|table (default markdown)
  OCR_PDF_ENGINE   auto|vision|text (default auto: per page — a substantial text
                   layer is used as-is, a page without one is vision-OCR'd, and a
                   text page that ALSO carries big raster images gets a HYBRID
                   text+vision pass so figures/diagrams aren't silently dropped)
  OCR_PREPROCESS   "1" to enable Ollama-OCR's cv2 binarization (default off — a raw
                   image usually reads better on a vision LLM)
  OCR_DOCLING      "1" to route non-image docs through Docling (must be installed)
"""

import os
import re
import hmac
import time
import shutil
import tempfile
import urllib.parse

import pymupdf
from fastapi import FastAPI, Request, HTTPException
from starlette.concurrency import run_in_threadpool

from ocr_processor import OCRProcessor

# ---- config from env (set once by extract.js spawnExtract) ------------------
KEY = os.environ.get("EXTRACT_API_KEY", "")
OCR_MODEL = os.environ.get("OCR_MODEL", "llama3.2-vision:11b")
OCR_OLLAMA_URL = os.environ.get("OCR_OLLAMA_URL", "http://127.0.0.1:11434/api/generate")
OCR_FORMAT = os.environ.get("OCR_FORMAT", "markdown")
OCR_PDF_ENGINE = os.environ.get("OCR_PDF_ENGINE", "auto")   # auto | vision | text
OCR_PREPROCESS = os.environ.get("OCR_PREPROCESS", "0") == "1"
OCR_DOCLING = os.environ.get("OCR_DOCLING", "0") == "1"
# Read timeout (seconds) for the Ollama call — a non-stream vision generate sends
# nothing until done, so this bounds total per-page generation before we give up and
# free the worker (a stalled Ollama would otherwise hold a threadpool token forever).
OCR_HTTP_TIMEOUT = int(os.environ.get("OCR_HTTP_TIMEOUT", "600"))

OCR = OCRProcessor(model_name=OCR_MODEL, base_url=OCR_OLLAMA_URL, request_timeout=(10, OCR_HTTP_TIMEOUT))

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp", ".gif"}
TEXT_EXTS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".log", ".rst"}
HTML_EXTS = {".html", ".htm"}
# Content-Type → extension fallback when X-Filename carries no usable suffix.
CT_EXT = {
    "image/png": ".png", "image/jpeg": ".jpg", "image/webp": ".webp",
    "image/tiff": ".tiff", "image/bmp": ".bmp", "image/gif": ".gif",
    "application/pdf": ".pdf",
    "text/plain": ".txt", "text/markdown": ".md", "text/csv": ".csv",
    "text/html": ".html",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
}

app = FastAPI(title="lol-extract", docs_url=None, redoc_url=None)

_docling_conv = None


def _docling_converter():
    """Lazily build (and cache) a Docling converter — importing docling pulls a
    heavy tree, so we only touch it when OCR_DOCLING is on and a doc arrives."""
    global _docling_conv
    if _docling_conv is None:
        from docling.document_converter import DocumentConverter
        _docling_conv = DocumentConverter()
    return _docling_conv


def _page(content, page, source, engine):
    return {"page_content": content or "", "metadata": {"page": page, "source": source, "engine": engine}}


def _ext(filename, content_type):
    ext = os.path.splitext(filename or "")[1].lower()
    if ext:
        return ext
    return CT_EXT.get((content_type or "").split(";")[0].strip().lower(), "")


def _read_text(path):
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


def _strip_html(s):
    s = re.sub(r"(?is)<(script|style).*?</\1>", " ", s)
    s = re.sub(r"(?s)<[^>]+>", " ", s)
    s = re.sub(r"[ \t]{2,}", " ", s)
    return s.strip()


def _extract_docx(path):
    import docx  # python-docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path):
    from pptx import Presentation  # python-pptx
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"# Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        lines.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def _extract_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _ocr_image(path):
    """Run a vision-model OCR pass. Ollama-OCR never raises — it returns an
    'Error processing image: …' string — so surface that as a 502 (Ollama down /
    model missing) rather than silently ingesting the error text."""
    res = OCR.process_image(path, format_type=OCR_FORMAT, preprocess=OCR_PREPROCESS)
    if isinstance(res, str) and res.startswith("Error processing image:"):
        raise HTTPException(status_code=502, detail=res)
    return res


# `auto` PDF routing: a page with fewer text-layer chars than this is treated as
# scanned (vision OCR); a page above it that still carries big raster images gets a
# HYBRID pass (text layer + vision) so figures/diagrams aren't silently dropped —
# the failure mode that plain "text layer wins" had on design docs/slides.
PDF_MIN_TEXT_CHARS = 32
PDF_HYBRID_IMG_COVERAGE = 0.2   # ≥20% of the page area under raster images → hybrid


def _image_coverage(page):
    """Fraction of the page area covered by raster images (0..1). Vector-drawn
    charts aren't images and don't count — force OCR_PDF_ENGINE=vision for those."""
    try:
        area = page.rect.get_area()
        if not area:
            return 0.0
        covered = sum((pymupdf.Rect(im["bbox"]) & page.rect).get_area()
                      for im in page.get_image_info())
        return min(1.0, covered / area)
    except Exception:
        return 0.0


def _render_page(page, workdir, i):
    pix = page.get_pixmap(matrix=pymupdf.Matrix(2, 2))  # ~144 DPI
    png = os.path.join(workdir, f"page{i}.png")
    pix.save(png)
    return png


def _extract_pdf(src, filename, workdir):
    """Per page: the embedded text layer when it's substantial (born-digital),
    vision OCR when it isn't (scanned), and BOTH when a real text layer coexists
    with large images (the text alone would drop the figures). OCR_PDF_ENGINE
    forces text-only or vision-only."""
    doc = pymupdf.open(src)
    out = []
    try:
        # Multi-page docs get a "[Page N]" header on every page's content: it makes
        # per-page coverage VISIBLE in OWUI's extracted-text preview, and it lets a
        # page-targeted question ("what's on page 5?") actually retrieve that page's
        # chunks — page numbers otherwise exist only in metadata RAG can't match on.
        multi = doc.page_count > 1
        label = (lambda i, c: f"[Page {i + 1}]\n{c}") if multi else (lambda i, c: c)
        for i in range(doc.page_count):
            page = doc[i]
            text = page.get_text("text").strip()
            if OCR_PDF_ENGINE == "text":
                mode = "text"
            elif OCR_PDF_ENGINE == "vision":
                mode = "vision"
            elif len(text) < PDF_MIN_TEXT_CHARS:   # auto: no real text layer → scanned
                mode = "vision"
            elif _image_coverage(page) >= PDF_HYBRID_IMG_COVERAGE:
                mode = "hybrid"
            else:
                mode = "text"
            if mode == "vision":
                out.append(_page(label(i, _ocr_image(_render_page(page, workdir, i))), i + 1, filename, "vision"))
            elif mode == "hybrid":
                # A vision failure here only loses the figure text — keep the page's
                # text layer rather than failing the whole document.
                try:
                    ocr = _ocr_image(_render_page(page, workdir, i))
                    out.append(_page(label(i, f"{text}\n\n[Page image content]\n{ocr}"), i + 1, filename, "text+vision"))
                except HTTPException as e:
                    print(f"page {i + 1}: hybrid vision pass failed ({e.detail}); kept the text layer", flush=True)
                    out.append(_page(label(i, text), i + 1, filename, "text"))
            else:
                out.append(_page(label(i, text), i + 1, filename, "text"))
    finally:
        doc.close()
    return out or [_page("", 1, filename, "text")]


def _extract(body, filename, content_type):
    ext = _ext(filename, content_type)
    workdir = tempfile.mkdtemp(prefix="lolocr_")
    src = os.path.join(workdir, "input" + (ext or ".bin"))
    with open(src, "wb") as f:
        f.write(body)
    try:
        # Images always go to the vision model — that's the point of OCR here.
        if ext in IMAGE_EXTS:
            return [_page(_ocr_image(src), 1, filename, "vision")]
        # Everything else: Docling (if enabled) for the best office/table fidelity,
        # else our light extractors. Docling failure falls through to the light path.
        if OCR_DOCLING:
            try:
                conv = _docling_converter()
                md = conv.convert(src).document.export_to_markdown()
                return [_page(md, 1, filename, "docling")]
            except HTTPException:
                raise
            except Exception as e:  # noqa: BLE001 — degrade gracefully to light path
                print(f"[extract] docling failed ({e}); falling back to light extraction", flush=True)
        if ext == ".pdf":
            return _extract_pdf(src, filename, workdir)
        if ext == ".docx":
            return [_page(_extract_docx(src), 1, filename, "text")]
        if ext == ".pptx":
            return [_page(_extract_pptx(src), 1, filename, "text")]
        if ext == ".xlsx":
            return [_page(_extract_xlsx(src), 1, filename, "text")]
        if ext in TEXT_EXTS:
            return [_page(_read_text(src), 1, filename, "text")]
        if ext in HTML_EXTS:
            return [_page(_strip_html(_read_text(src)), 1, filename, "text")]
        # Legacy binary office (.doc/.ppt/.xls) + odt/epub/rtf aren't in the light
        # path; enabling ocr.docling routes them through Docling above instead of 415.
        raise HTTPException(
            status_code=415,
            detail=f'Unsupported file type "{ext or content_type}". Images + PDFs are OCR\'d; '
                   f"docx/pptx/xlsx/text extract natively; enable ocr.docling for other office formats.",
        )
    finally:
        shutil.rmtree(workdir, ignore_errors=True)


def _check_auth(authorization):
    if not KEY:
        raise HTTPException(status_code=500, detail="OCR service misconfigured: no API key set.")
    if not hmac.compare_digest(authorization or "", f"Bearer {KEY}"):
        raise HTTPException(status_code=401, detail="Unauthorized")


@app.get("/health")
async def health():
    return {"status": "ok", "model": OCR_MODEL, "docling": OCR_DOCLING}


@app.put("/process")
async def process(request: Request):
    _check_auth(request.headers.get("authorization", ""))
    body = await request.body()
    if not body:
        raise HTTPException(status_code=400, detail="Empty request body.")
    filename = urllib.parse.unquote(request.headers.get("x-filename", "") or "upload")
    content_type = request.headers.get("content-type", "")
    # OCR / extraction is blocking (HTTP to Ollama, cv2, pymupdf) — run it off the
    # event loop so concurrent uploads don't serialize on the single async worker.
    t0 = time.monotonic()
    pages = await run_in_threadpool(_extract, body, filename, content_type)
    # One summary line per document in the farm log ([extract] prefix added by the
    # supervisor) — the operator's proof that OWUI is routed here at all, and the
    # first thing to read when an extraction "missed" content: it shows how each
    # page was handled (text layer / vision OCR / hybrid).
    engines = {}
    chars = 0
    for p in pages:
        engines[p["metadata"].get("engine", "?")] = engines.get(p["metadata"].get("engine", "?"), 0) + 1
        chars += len(p.get("page_content") or "")
    summary = " + ".join(f"{n} {e}" for e, n in sorted(engines.items()))
    print(f"{filename}: {len(pages)} page(s) → {summary} · {chars} chars · {time.monotonic() - t0:.1f}s", flush=True)
    return pages


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=int(os.environ.get("EXTRACT_PORT", "8890")))
